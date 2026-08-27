#!/usr/bin/env python3
"""img2pptx MCP server (stdio).

Exposes the deterministic core of the img2pptx skill as MCP tools so any
MCP-capable agent (Google Antigravity, Claude Desktop, Cursor, ...) can
render/compare/package without writing ad-hoc scripts.

Run:
    python3 mcp/img2pptx_mcp/server.py          # stdio MCP server
Install (Antigravity): see docs/antigravity.md or scripts/sync_antigravity.py
"""
from __future__ import annotations

import hashlib
import io
import json
import os
import re
import shutil
import sys
import zipfile
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[2]
FONTS_CANDIDATES = [
    Path(os.environ.get("IMG2PPTX_FONTS_DIR", "/nonexistent")),
    REPO_ROOT / "task" / "fonts",
    REPO_ROOT / "mcp" / "img2pptx_mcp" / "fonts",
    Path.home() / ".local" / "share" / "img2pptx" / "fonts",
]


def _fonts() -> list[str]:
    out = []
    for d in FONTS_CANDIDATES:
        if d.is_dir():
            out += [str(p) for p in sorted(d.glob("*.ttf"))]
    # metric-compatible fallbacks first, then whatever the OS has
    out.append("")  # let resvg use system fonts too
    return out


def _deps() -> dict:
    status = {}
    for mod, pkg in [("PIL", "pillow"), ("numpy", "numpy"),
                     ("pptx", "python-pptx"), ("resvg_py", "resvg-py")]:
        try:
            __import__(mod)
            status[pkg] = "ok"
        except Exception as e:  # noqa: BLE001
            status[pkg] = f"missing ({e.__class__.__name__})"
    status["fonts_found"] = len([f for f in _fonts() if f])
    return status


try:  # mcp v1.x
    from mcp.server.fastmcp import FastMCP as MCPServer
except ModuleNotFoundError:  # mcp 2.x: FastMCP renamed to MCPServer
    from mcp.server.mcpserver import MCPServer  # noqa: N813

mcp = MCPServer("img2pptx")


@mcp.tool()
def check_environment() -> str:
    """Verify the render/compare/pptx toolchain of the img2pptx skill.

    Returns JSON with dependency status, font count and repo root.
    """
    return json.dumps({
        "dependencies": _deps(),
        "repo_root": str(REPO_ROOT),
        "skill_md": str(REPO_ROOT / "skills" / "img2pptx" / "SKILL.md"),
        "python": sys.version.split()[0],
    }, ensure_ascii=False)


@mcp.tool()
def skill_info() -> str:
    """Return the img2pptx workflow digest (protocol steps + deliverables).

    Read this before reconstructing an image so the pipeline order and the
    mandatory QA artifacts are known without opening the full SKILL.md.
    """
    steps = [
        "1 inspect input + normalize to PNG (keep original)",
        "2 model semantic units / combination submodules / atomic elements",
        "3 extract semantic constraints (source_observations, invariants, negatives)",
        "4 build component_manifest.json (bbox, hierarchy, audit map)",
        "5 layout skeleton audit (panels, alignment, arrows)",
        "6 draw modules standalone (border layering: fill bg + top outline)",
        "7 standalone integrity audit (clipping, dependencies)",
        "8 assemble full.svg with nested <g id> hierarchy",
        "9 containment + alignment + border-layering audits",
        "10 semantic constraint + coverage audits (hard gate)",
        "11 visual similarity audit (MAE whole + per component, overlay, diff)",
        "12 build one-slide PPTX (SVG primary via svgBlip, PNG fallback)",
        "13 pptx package audit (hash match, rels, content types, one slide)",
        "14 mandatory correction loop; bounded MAE optimization (<=3 attempts)",
        "15 deliver final.pptx + full.svg + modules + manifest + qa/*",
    ]
    return json.dumps({
        "skill": "img2pptx",
        "protocol": steps,
        "deterministic_tools": [t for t in ("normalize_image", "render_svg",
                                            "compare_images", "embed_svg_pptx",
                                            "audit_pptx")],
        "rules": [
            "SVG is the primary representation; never ship a full-slide bitmap only",
            "text stays editable <text>; photos may be minimal raster crops only",
            "semantic errors override low MAE; never fabricate illegible values",
        ],
        "full_spec": str(REPO_ROOT / "skills" / "img2pptx" / "SKILL.md"),
    }, ensure_ascii=False)


@mcp.tool()
def normalize_image(image_path: str, out_png: str) -> str:
    """Normalize any raster input to an upright RGB PNG (original untouched).

    Applies EXIF orientation, converts to RGB, saves lossless PNG.
    Returns JSON with source/normalized paths and dimensions.
    """
    from PIL import Image
    src = Path(image_path).expanduser().resolve()
    im = Image.open(src)
    try:
        im = im.convert("RGB")
    except Exception:
        pass
    try:
        from PIL import ImageOps
        im = ImageOps.exif_transpose(im)
    except Exception:
        pass
    im = im.convert("RGB")
    out = Path(out_png).expanduser().resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    im.save(out, "PNG", optimize=True)
    return json.dumps({
        "source": str(src), "normalized": str(out),
        "size": im.size, "aspect": round(im.width / im.height, 4),
    }, ensure_ascii=False)


@mcp.tool()
def render_svg(svg_path: str, out_png: str, width: int | None = None,
               height: int | None = None) -> str:
    """Render an SVG to PNG with resvg (white background, metric-compatible
    font fallbacks Caladea/Carlito when bundled fonts are present).

    Returns JSON with the output path and pixel size.
    """
    import resvg_py
    from PIL import Image
    p = Path(svg_path).expanduser().resolve()
    kw = {"background": "#ffffff"}
    fonts = [f for f in _fonts() if f]
    if fonts:
        kw.update(font_files=fonts, serif_family="Caladea",
                  sans_serif_family="Carlito")
    if width:
        kw["width"] = width
    if height:
        kw["height"] = height
    png = resvg_py.svg_to_bytes(p.read_text(), **kw)
    im = Image.open(io.BytesIO(png)).convert("RGB")
    out = Path(out_png).expanduser().resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    im.save(out, "PNG")
    return json.dumps({"svg": str(p), "png": str(out), "size": im.size},
                      ensure_ascii=False)


@mcp.tool()
def compare_images(source_path: str, render_path: str,
                   out_prefix: str | None = None, grid: int = 6) -> str:
    """Compare a reconstruction render against its normalized source.

    Renders are resized to the source size first. Returns JSON with
    whole-image MAE (0-255), missing-ink and extra-ink fractions, a grid
    MAE heatmap and per-grid strong-error rates. When out_prefix is given,
    writes <prefix>_overlay.png and <prefix>_diff.png (3x amplified).
    """
    import numpy as np
    from PIL import Image
    a = np.asarray(Image.open(source_path).convert("RGB"), float)
    im_b = Image.open(render_path).convert("RGB")
    if im_b.size != (a.shape[1], a.shape[0]):
        im_b = im_b.resize((a.shape[1], a.shape[0]), Image.LANCZOS)
    b = np.asarray(im_b, float)
    d = np.abs(a - b).mean(axis=2)
    ga, gb = a.mean(axis=2), b.mean(axis=2)
    res = {
        "source": source_path, "render": render_path,
        "size": [a.shape[1], a.shape[0]],
        "mae": round(float(d.mean()), 2),
        "missing_ink_pct": round(float(((ga < 225) & (gb > 245)).mean() * 100), 2),
        "extra_ink_pct": round(float(((ga > 245) & (gb < 225)).mean() * 100), 2),
        "strong_err_pct": round(float((d > 100).mean() * 100), 2),
    }
    g = max(1, min(12, grid))
    h, w = d.shape
    res["mae_grid"] = [[round(float(d[r * h // g:(r + 1) * h // g,
                                      c * w // g:(c + 1) * w // g].mean()), 1)
                        for c in range(g)] for r in range(g)]
    if out_prefix:
        pre = Path(out_prefix).expanduser().resolve()
        pre.parent.mkdir(parents=True, exist_ok=True)
        Image.blend(Image.fromarray(a.astype("uint8")),
                    Image.fromarray(b.astype("uint8")), 0.5).save(
            str(pre) + "_overlay.png")
        Image.fromarray((np.clip(d * 3, 0, 255)).astype("uint8")).save(
            str(pre) + "_diff.png")
        res["artifacts"] = [str(pre) + "_overlay.png", str(pre) + "_diff.png"]
    return json.dumps(res, ensure_ascii=False)


def _svgblip_pptx(svg_path: Path, png_path: Path, out_path: Path,
                  width_in: float | None, height_in: float | None) -> dict:
    """Build a one-slide PPTX whose picture carries an SVG blip extension."""
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image

    png = Image.open(png_path)
    aspect = png.width / png.height
    if width_in and height_in:
        w, h = width_in, height_in
    elif width_in:
        w, h = width_in, width_in / aspect
    elif height_in:
        w, h = height_in * aspect, height_in
    else:
        w, h = 13.333, 13.333 / aspect

    prs = Presentation()
    prs.slide_width = Inches(w)
    prs.slide_height = Inches(h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(png_path), 0, 0,
                             width=prs.slide_width, height=prs.slide_height)
    tmp = out_path.with_suffix(".raw.pptx")
    prs.save(tmp)

    svg_bytes = svg_path.read_bytes()
    zin = zipfile.ZipFile(tmp, "r")
    zout = zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED)
    rels = zin.read("ppt/slides/_rels/slide1.xml.rels").decode()
    m = re.search(r'Id="(rId\d+)"[^>]*Target="\.\./media/([^"]+\.png)"', rels)
    if not m:
        raise ValueError("picture relationship not found in slide1 rels")
    png_rid = m.group(1)
    svg_rid = "rIdSvg1"
    rels_new = rels.replace(
        "</Relationships>",
        f'<Relationship Id="{svg_rid}" Type="http://schemas.openxmlformats.org/'
        f'officeDocument/2006/relationships/image" '
        f'Target="../media/image_svg.svg"/></Relationships>')
    slide_xml = zin.read("ppt/slides/slide1.xml").decode()
    blip = f'<a:blip r:embed="{png_rid}"/>'
    if blip not in slide_xml:
        raise ValueError("blip element not found in slide1.xml")
    slide_xml = slide_xml.replace(
        blip,
        f'<a:blip r:embed="{png_rid}">'
        f'<a:extLst><a:ext uri="{{96DAC541-7B7A-43D3-8B79-37D633B846F1}}">'
        f'<asvg:svgBlip xmlns:asvg="http://schemas.microsoft.com/office/'
        f'drawing/2016/SVG/main" r:embed="{svg_rid}"/></a:ext></a:extLst>'
        f'</a:blip>')
    ct = zin.read("[Content_Types].xml").decode()
    if "svg" not in ct:
        ct = ct.replace("</Types>",
                        '<Default Extension="svg" ContentType="image/svg+xml"/>'
                        "</Types>")
    for item in zin.infolist():
        data = zin.read(item.filename)
        if item.filename == "ppt/slides/_rels/slide1.xml.rels":
            data = rels_new.encode()
        elif item.filename == "ppt/slides/slide1.xml":
            data = slide_xml.encode()
        elif item.filename == "[Content_Types].xml":
            data = ct.encode()
        zout.writestr(item, data)
    zout.writestr("ppt/media/image_svg.svg", svg_bytes)
    zout.close()
    zin.close()
    tmp.unlink(missing_ok=True)
    return {"slide_in": [round(w, 3), round(h, 3)], "aspect": round(aspect, 4)}


@mcp.tool()
def embed_svg_pptx(svg_path: str, preview_png: str, out_pptx: str,
                   width_in: float | None = None,
                   height_in: float | None = None) -> str:
    """Package full.svg (+ its PNG render) into a one-slide editable PPTX.

    The picture carries the SVG as primary representation (svgBlip) with the
    PNG as fallback; the package is then audited. Returns JSON audit results.
    """
    svg = Path(svg_path).expanduser().resolve()
    png = Path(preview_png).expanduser().resolve()
    out = Path(out_pptx).expanduser().resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    info = _svgblip_pptx(svg, png, out, width_in, height_in)
    audit = json.loads(audit_pptx(str(out)))
    audit.update(info)
    return json.dumps(audit, ensure_ascii=False)


@mcp.tool()
def audit_pptx(pptx_path: str) -> str:
    """Audit a PPTX package: SVG part present + hash, svgBlip wiring, rels,
    content types, PNG fallback, slide count, group/text counts in the
    embedded SVG, unsupported-feature scan. Returns JSON with passed flag.
    """
    p = Path(pptx_path).expanduser().resolve()
    z = zipfile.ZipFile(p)
    names = z.namelist()
    checks = []
    svg_parts = [n for n in names if n.endswith(".svg")]
    checks.append({"check": "svg_part_present", "parts": svg_parts,
                   "passed": bool(svg_parts)})
    h = None
    if svg_parts:
        emb = z.read(svg_parts[0])
        h = hashlib.sha256(emb).hexdigest()
        s = emb.decode("utf-8", "replace")
        checks.append({
            "check": "embedded_svg_integrity",
            "sha256": h,
            "n_groups": s.count("<g "),
            "n_texts": s.count("<text"),
            "n_data_uri_images": s.count("data:image/"),
            "passed": True})
        sx = ""
        for n in names:
            if re.fullmatch(r"ppt/slides/slide\d+\.xml", n):
                sx += z.read(n).decode("utf-8", "replace")
        checks.append({"check": "svgblip_extension_in_blip",
                       "passed": "svgBlip" in sx})
        unsupported = [f for f in ("foreignObject", "<filter")
                       if f in s]
        checks.append({"check": "no_unsupported_features",
                       "found": unsupported, "passed": not unsupported})
    rels_svg = any(re.search(rb'Id="[^"]+"[^>]*Target="[^"]*\.svg"',
                             z.read(n)) for n in names if n.endswith(".rels"))
    checks.append({"check": "rels_reference_svg", "passed": rels_svg})
    png_fallback = any(n.endswith(".png") for n in names if
                       n.startswith("ppt/media/"))
    checks.append({"check": "png_fallback_present", "passed": png_fallback})
    n_slides = len([n for n in names if re.fullmatch(r"ppt/slides/slide\d+\.xml",
                                                     n)])
    checks.append({"check": "exactly_one_slide", "value": n_slides,
                   "passed": n_slides == 1})
    ctx = z.read("[Content_Types].xml").decode()
    checks.append({"check": "svg_content_type_registered",
                   "passed": "image/svg+xml" in ctx or not svg_parts})
    return json.dumps({"pptx": str(p), "bytes": p.stat().st_size,
                       "checks": checks,
                       "passed": all(c["passed"] for c in checks)},
                      ensure_ascii=False)


if __name__ == "__main__":
    mcp.run()
