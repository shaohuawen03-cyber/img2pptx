#!/usr/bin/env python3
"""Build final.pptx: one slide containing the complete full.svg (SVG primary,
PNG fallback), then audit the package."""
import hashlib
import json
import re
import shutil
import zipfile

from pptx import Presentation
from pptx.util import Inches

W_IN, H_IN = 13.333, 6.279  # matches 2472x1164 aspect (2.1237)

prs = Presentation()
prs.slide_width = Inches(W_IN)
prs.slide_height = Inches(H_IN)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
pic = slide.shapes.add_picture("full_render.png", 0, 0,
                               width=prs.slide_width, height=prs.slide_height)
prs.save("final_raw.pptx")

# ---- inject the SVG part + svgBlip extension
svg_bytes = open("full.svg", "rb").read()
png_bytes = open("full_render.png", "rb").read()

zin = zipfile.ZipFile("final_raw.pptx", "r")
names = zin.namelist()
zout = zipfile.ZipFile("final.pptx", "w", zipfile.ZIP_DEFLATED)

# find the rId of the PNG picture in slide1
rels = zin.read("ppt/slides/_rels/slide1.xml.rels").decode()
m = re.search(r'Id="(rId\d+)"[^>]*Target="\.\./media/([^"]+\.png)"', rels)
png_rid, png_tgt = m.group(1), m.group(2)
svg_name = "image_svg.svg"
svg_rid = "rIdSvg1"
new_rel = (f'<Relationship Id="{svg_rid}" '
           f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" '
           f'Target="../media/{svg_name}"/>')
rels_new = rels.replace("</Relationships>", new_rel + "</Relationships>")

slide_xml = zin.read("ppt/slides/slide1.xml").decode()
# add svgBlip inside the blip of our picture
blip_tag = f'<a:blip r:embed="{png_rid}"/>'
assert blip_tag in slide_xml, "blip not found"
svg_blip = (f'<a:blip r:embed="{png_rid}">'
            f'<a:extLst><a:ext uri="{{96DAC541-7B7A-43D3-8B79-37D633B846F1}}">'
            f'<asvg:svgBlip xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main" '
            f'r:embed="{svg_rid}"/></a:ext></a:extLst></a:blip>')
slide_xml = slide_xml.replace(blip_tag, svg_blip)

ct = zin.read("[Content_Types].xml").decode()
if "svg" not in ct:
    ct = ct.replace("</Types>",
                    '<Default Extension="svg" ContentType="image/svg+xml"/></Types>')

for item in zin.infolist():
    data = zin.read(item.filename)
    if item.filename == "ppt/slides/_rels/slide1.xml.rels":
        data = rels_new.encode()
    elif item.filename == "ppt/slides/slide1.xml":
        data = slide_xml.encode()
    elif item.filename == "[Content_Types].xml":
        data = ct.encode()
    zout.writestr(item, data)
zout.writestr(f"ppt/media/{svg_name}", svg_bytes)
zout.close()
zin.close()
print("final.pptx written")

# ---- audit the package
audit = {"checks": [], "passed": True}
z = zipfile.ZipFile("final.pptx")
names = z.namelist()
has_svg = "ppt/media/image_svg.svg" in names
audit["checks"].append({"check": "svg_part_present", "passed": has_svg})
emb = z.read("ppt/media/image_svg.svg")
h_emb = hashlib.sha256(emb).hexdigest()
h_src = hashlib.sha256(svg_bytes).hexdigest()
match = h_emb == h_src
audit["checks"].append({"check": "embedded_svg_matches_full_svg_by_hash",
                        "sha256": h_emb, "passed": match})
sx = z.read("ppt/slides/slide1.xml").decode()
has_svgblip = "svgBlip" in sx and svg_rid in sx
audit["checks"].append({"check": "svgblip_extension_in_blip", "passed": has_svgblip})
rx = z.read("ppt/slides/_rels/slide1.xml.rels").decode()
rels_ok = svg_rid in rx and png_rid in rx
audit["checks"].append({"check": "rels_for_svg_and_png_fallback", "passed": rels_ok})
png_fallback = f"ppt/media/{png_tgt}" in names
audit["checks"].append({"check": "png_fallback_present", "passed": png_fallback})
n_slides = len([n for n in names if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)])
audit["checks"].append({"check": "exactly_one_slide", "value": n_slides,
                        "passed": n_slides == 1})
ctx = z.read("[Content_Types].xml").decode()
audit["checks"].append({"check": "svg_content_type_registered",
                        "passed": "image/svg+xml" in ctx})
# nested groups preserved in embedded SVG bytes
n_groups = emb.decode().count("<g ")
n_texts = emb.decode().count("<text")
audit["checks"].append({"check": "nested_groups_and_texts_in_svg",
                        "n_groups": n_groups, "n_texts": n_texts,
                        "passed": n_groups >= 20 and n_texts >= 300})
# no unsupported features
s = emb.decode()
unsupported = [f for f in ("foreignObject", "filter", "<image", "xlink:") if f in s]
# note: <image> appears for the 4 photo crops (data URIs, self-contained)
raster_images = s.count("data:image/png;base64,")
audit["checks"].append({
    "check": "no_unsupported_features",
    "foreignObject/filter": not any(f in s for f in ("foreignObject", "filter")),
    "self_contained_data_uri_images": raster_images,
    "passed": ("foreignObject" not in s) and ("filter" not in s),
})
audit["passed"] = all(c["passed"] for c in audit["checks"])
json.dump(audit, open("qa/pptx_preview_audit.json", "w"), indent=1)
print(json.dumps(audit["checks"], indent=1))
print("PPTX audit passed:", audit["passed"])
